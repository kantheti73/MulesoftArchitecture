"""
Generates presentations/MuleSoft-APIGateway-OnPrem.pptx

A 10-slide deck for the audience that has already decided on-prem (SaaS
isn't viable). Companion to build_pptx.py (which covers the SaaS path).

Run:
    pip install --user python-pptx
    # if the diagram source changed:
    python presentations/diagrams/onprem_architecture.py
    # always:
    python presentations/build_onprem_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# -------- Branding / palette --------
NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xDE, 0xEB, 0xF7)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
GREEN_OK = RGBColor(0x33, 0x99, 0x33)
RED_NO = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_PATH = Path("presentations/MuleSoft-APIGateway-OnPrem.pptx")
ARCH_IMG = Path("presentations/img/onprem-architecture.png")


# -------- Helpers (duplicated from build_pptx.py so the two decks stay independent) --------
def add_solid_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = [text] if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.alignment = align
    return tb


def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=DARK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = "•  " + text
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
    return tb


def add_header(slide, title_text, subtitle_text=None):
    add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), NAVY)
    add_text(slide, title_text, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.55),
             font_size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle_text:
        add_text(slide, subtitle_text, Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.35),
                 font_size=13, color=BLUE, anchor=MSO_ANCHOR.TOP)


def add_footer(slide, slide_no, total):
    add_solid_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), LIGHT_GRAY)
    add_text(slide, "MuleSoft API Gateway — On-Prem Deployment",
             Inches(0.4), Inches(7.22), Inches(8), Inches(0.28),
             font_size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, f"{slide_no} / {total}",
             Inches(12.0), Inches(7.22), Inches(1.0), Inches(0.28),
             font_size=10, color=DARK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def make_table(slide, data, left, top, width, height,
               header_row=True, first_col_bold=False,
               header_bg=NAVY, header_fg=WHITE,
               body_bg_alt=LIGHT_GRAY, body_fg=DARK,
               header_font_size=12, body_font_size=11,
               col_widths=None):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    if col_widths is not None:
        total_emu = sum(col_widths)
        scale = width / total_emu if total_emu else 1.0
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(w * scale)

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

            is_header = (header_row and r == 0)
            bg = header_bg if is_header else (body_bg_alt if r % 2 == 0 else WHITE)
            fg = header_fg if is_header else body_fg
            font_size = header_font_size if is_header else body_font_size

            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.color.rgb = fg
            run.font.bold = is_header or (first_col_bold and c == 0 and not is_header)
            p.alignment = PP_ALIGN.LEFT
    return table


# -------- Slide builders --------
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_solid_rect(s, Inches(0), Inches(2.3), SLIDE_W, Inches(2.4), NAVY)
    add_solid_rect(s, Inches(0), Inches(4.7), SLIDE_W, Inches(0.12), ORANGE)

    add_text(s, "MuleSoft API Gateway", Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.9),
             font_size=44, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "On-Prem Deployment", Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
             font_size=26, color=LIGHT_BLUE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Anypoint Flex Gateway in your data center — active/active across DCs with full operational control",
             Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.7),
             font_size=16, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Prepared by the Integration Architecture team",
             Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             font_size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "v1.0  ·  Companion to the SaaS deck", Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
             font_size=11, color=BLUE, anchor=MSO_ANCHOR.MIDDLE)
    return s


def slide_why_onprem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Why on-prem",
               "The scenarios where SaaS isn't viable, and we accept the operational responsibility in exchange for full control")

    bullets = [
        "Regulatory regime mandates on-prem — some sovereign-cloud, defense, classified, or healthcare contexts where any SaaS-managed runtime is non-starter.",
        "Risk appetite cannot accept ANY MuleSoft personnel access to runtime infrastructure — even contractually constrained (SOC 2 + DPA insufficient).",
        "We already operate a hardened on-prem PII platform (HSM, dedicated SOC, DLP, SIEM, internal CA) — adding the gateway is incremental, not greenfield.",
        "Sustained volume > 5–10M calls/day — at this scale, amortized on-prem hardware undercuts the SaaS license model.",
        "True air-gap operation required — no outbound connection to anypoint.mulesoft.com is permitted (forces Local mode; no SaaS control plane at all).",
        "Strategic preference to control upgrade cadence, monitoring stack, incident response, and audit trails entirely in-house.",
    ]
    add_bullets(s, bullets, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.8),
                font_size=15, color=DARK)

    # Trade-off callout
    add_solid_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7), LIGHT_BLUE)
    add_text(s, "The trade-off you accept",
             Inches(0.7), Inches(6.42), Inches(4), Inches(0.32),
             font_size=13, bold=True, color=NAVY)
    add_text(s, "OS patching, runtime upgrades, HA, scaling, monitoring, incident response — all yours. The SaaS path absorbs these; on-prem doesn't.",
             Inches(0.7), Inches(6.72), Inches(12.0), Inches(0.36),
             font_size=11, color=DARK)
    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "On-prem architecture overview",
               "Active/active two-DC topology with F5 GSLB, GitOps-driven config, Vault-managed secrets")

    if ARCH_IMG.exists():
        s.shapes.add_picture(str(ARCH_IMG), Inches(0.3), Inches(1.5),
                             width=Inches(8.9), height=Inches(5.4))
    else:
        add_text(s, "[on-prem architecture image not generated]", Inches(0.5), Inches(2),
                 Inches(8.5), Inches(2), font_size=14, color=RED_NO)

    add_solid_rect(s, Inches(9.4), Inches(1.5), Inches(3.7), Inches(5.4), LIGHT_BLUE)
    add_text(s, "Highlights", Inches(9.55), Inches(1.6), Inches(3.5), Inches(0.4),
             font_size=16, bold=True, color=NAVY)
    highlights = [
        "Active/active across two DCs — both serve traffic simultaneously",
        "F5 BIG-IP DNS / GSLB fronts both DCs with 30s TTL health-checked failover",
        "Per-DC F5 LTM in front of 2 Flex Gateway replicas (4 replicas total)",
        "Git is the config source of truth → Ansible AWX pushes to both DCs in parallel",
        "HashiCorp Vault performance replication keeps secrets in-sync sub-second",
        "Optional outbound mTLS to Anypoint Control Plane (Connected mode only)",
        "MS Integration Stack replicated across DCs (own DR mechanism)",
        "Splunk HEC + Datadog APM dual-shipped from each replica",
    ]
    add_bullets(s, highlights, Inches(9.5), Inches(2.05), Inches(3.55), Inches(4.8),
                font_size=10, color=DARK, line_spacing=1.10)
    return s


def slide_modes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Deployment modes — Connected vs Local",
               "Same Flex Gateway binary; the difference is whether you stay tied to the Anypoint Control Plane")

    data = [
        ["Dimension", "Connected mode (recommended)", "Local mode (fully air-gapped)"],
        ["Outbound connection to anypoint.mulesoft.com", "Required (HTTPS 443, mTLS)", "None — fully isolated"],
        ["Source of truth for policy config", "Anypoint API Manager (SaaS)", "YAML files in your Git repo"],
        ["Policy push mechanism", "Long-polled from Anypoint, ~2–10s end-to-end", "CI pipeline → Ansible AWX → SIGHUP hot reload"],
        ["Central API catalog (Exchange)", "Yes — single source", "No — your own catalog (Backstage / wiki)"],
        ["Anypoint Monitoring", "Yes — built-in dashboards + analytics", "Not available; SIEM-only"],
        ["Audit log of admin actions", "Anypoint Audit Log + forwarded to SIEM", "SIEM only (Git history + Ansible logs)"],
        ["Resilience to anypoint.mulesoft.com outage", "Local policy cache survives ~24h", "N/A — no dependency"],
        ["Licensing", "Same Anypoint subscription", "Same Anypoint subscription"],
        ["Best for", "Default on-prem path; SaaS control plane fine to use", "Hard air-gap mandates / strict no-egress regimes"],
    ]
    make_table(s, data, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.45),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=11,
               col_widths=[Inches(3.3), Inches(4.6), Inches(4.6)])
    return s


def slide_vm_vs_k8s(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Install path — Standalone VM vs Kubernetes",
               "Both supported; pick by existing platform footprint and team skills")

    data = [
        ["Aspect", "Standalone VM (RHEL / Ubuntu)", "Kubernetes (OpenShift / Rancher / EKS-A)"],
        ["Install mechanism", "RPM / DEB package + systemd", "Helm chart + Deployment + Service"],
        ["HA model", "Manual — F5 LTM with VRRP / keepalived", "Native — replicas + PodDisruptionBudget + anti-affinity"],
        ["Self-healing on crash", "systemd Restart=on-failure", "Liveness/readiness probes + kubelet"],
        ["Scaling out", "Image + provision new VM + add to LB pool", "Kubectl scale OR HPA in seconds"],
        ["Rolling upgrades", "Manual orchestration (Ansible Tower)", "Native Deployment rolling strategy"],
        ["Operational tooling", "journalctl, systemctl, F5 console — familiar", "kubectl, Helm, K8s dashboards — specialist"],
        ["Compute overhead per node", "~50 MB (just Flex)", "~50 MB Flex + 0.5 vCPU + 1 GB K8s control-plane tax per worker"],
        ["Config drift control", "Ansible (manual discipline)", "GitOps + ArgoCD / Flux — declarative"],
        ["Cluster upgrade complexity", "N/A", "Quarterly project; rollback plan required"],
        ["Best fit", "✔  Your team already runs Linux + F5; no existing K8s footprint", "✔  K8s already a platform standard; cloud-native bet"],
    ]
    make_table(s, data, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.45),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=10,
               col_widths=[Inches(2.8), Inches(4.85), Inches(4.85)])
    return s


def slide_sizing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Hardware sizing — 5M calls/day target",
               "Provision for peak, monitor utilization, scale per the curve")

    # Left — per-replica spec
    add_text(s, "Per-replica spec (production)", Inches(0.5), Inches(1.55), Inches(6), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    perspec = [
        ["Resource", "Value"],
        ["vCPU", "4 (2 minimum)"],
        ["RAM", "8 GB (2 GB minimum)"],
        ["Disk", "100 GB SSD (~10 GB/day logs at 5M)"],
        ["NIC", "1 Gbps sufficient; 10 Gbps DC fabric"],
        ["OS", "RHEL 9 / Ubuntu 22.04 LTS"],
        ["Kernel", "5.14+ (better TCP perf)"],
        ["Time sync", "NTP, max 60s skew"],
    ]
    make_table(s, perspec, Inches(0.5), Inches(1.95), Inches(6.0), Inches(2.5),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=11,
               col_widths=[Inches(2.5), Inches(3.5)])

    # Right — 5M/day Prod target config
    add_text(s, "Target for 5M calls/day (Prod)", Inches(6.8), Inches(1.55), Inches(6), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    target = [
        ["Setting", "Value"],
        ["Peak TPS (5×) / spike (10×)", "~600 / ~1,200"],
        ["Replicas per DC", "4"],
        ["Number of DCs", "2 (active/active)"],
        ["Total replicas", "8"],
        ["Per-replica", "4 vCPU / 8 GB RAM"],
        ["Total Prod compute footprint", "32 vCPU + 64 GB RAM"],
        ["Headroom vs 1,200 TPS spike", "~2× (monitor closely)"],
    ]
    make_table(s, target, Inches(6.8), Inches(1.95), Inches(6.2), Inches(2.5),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=11,
               col_widths=[Inches(3.0), Inches(3.2)])

    # Bottom — sizing curve
    add_text(s, "Sizing curve — when to upsize",
             Inches(0.5), Inches(4.65), Inches(6), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    curve = [
        ["Volume", "Replicas / DC", "Per-replica spec", "Total compute (both DCs)"],
        ["≤ 100K/day", "2", "2 vCPU / 4 GB", "8 vCPU + 16 GB"],
        ["1M/day", "2", "4 vCPU / 8 GB", "16 vCPU + 32 GB"],
        ["5M/day (current target)", "4", "4 vCPU / 8 GB", "32 vCPU + 64 GB"],
        ["10M/day", "4", "8 vCPU / 16 GB", "64 vCPU + 128 GB"],
        ["25M/day", "6", "8 vCPU / 16 GB", "96 vCPU + 192 GB"],
        ["50M/day", "8+", "8 vCPU / 16 GB + kernel tuning", "128+ vCPU + 256+ GB"],
    ]
    make_table(s, curve, Inches(0.5), Inches(5.05), Inches(12.5), Inches(2.0),
               header_bg=BLUE, header_fg=WHITE, header_font_size=12, body_font_size=10,
               col_widths=[Inches(3.0), Inches(2.5), Inches(4.0), Inches(3.0)])
    return s


def slide_env_matrix(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Environments — DEV / QA / Acceptance / Prod",
               "QA mirrors Prod for accurate load testing — that decision drives roughly half of gateway compute spend")

    # Per-env sizing table
    data = [
        ["Setting", "DEV", "QA (= Prod)", "Acceptance / UAT", "PROD"],
        ["Volume target", "sandbox", "full 5M/day load test", "smoke + integration", "5M/day"],
        ["DC topology", "single DC", "2 DCs active/active", "2 DCs active/active", "2 DCs active/active"],
        ["Replicas per DC", "1", "4", "2", "4"],
        ["Total replicas", "1", "8", "4", "8"],
        ["Per-replica spec", "2 vCPU / 2 GB", "4 vCPU / 8 GB", "2 vCPU / 4 GB", "4 vCPU / 8 GB"],
        ["Compute total", "2 + 2", "32 + 64", "8 + 16", "32 + 64"],
        ["Redis", "1 dev node", "Sentinel 3-node × 2 DCs", "Sentinel 3-node × 1 DC", "Sentinel 3-node × 2 DCs"],
        ["Service Bus tier (if used)", "Standard", "Premium (load-test data)", "Standard or Premium", "Premium"],
        ["HA / DR drills validated here?", "No", "Yes — full Prod runbook", "Yes", "Yes (production)"],
    ]
    make_table(s, data, Inches(0.3), Inches(1.55), Inches(12.7), Inches(3.6),
               header_bg=NAVY, header_fg=WHITE, header_font_size=11, body_font_size=10,
               col_widths=[Inches(2.5), Inches(2.0), Inches(2.6), Inches(2.6), Inches(2.0)])

    # Aggregated footprint callout
    add_solid_rect(s, Inches(0.3), Inches(5.25), Inches(8.4), Inches(1.85), LIGHT_BLUE)
    add_text(s, "Aggregated footprint across all 4 environments",
             Inches(0.5), Inches(5.32), Inches(8.0), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    agg = [
        ["Resource", "DEV", "QA", "UAT", "PROD", "Total"],
        ["Gateway replicas", "1", "8", "4", "8", "21"],
        ["Gateway vCPU", "2", "32", "8", "32", "74"],
        ["Gateway RAM (GB)", "2", "64", "16", "64", "146"],
        ["Redis nodes", "1", "6", "3", "6", "16"],
        ["TOTAL compute (vCPU + GB)", "4 + 4", "44 + 76", "14 + 22", "44 + 76", "~106 vCPU + ~178 GB"],
    ]
    make_table(s, agg, Inches(0.5), Inches(5.68), Inches(8.0), Inches(1.4),
               header_bg=BLUE, header_fg=WHITE, header_font_size=10, body_font_size=9,
               col_widths=[Inches(2.4), Inches(0.8), Inches(0.8), Inches(0.8), Inches(0.8), Inches(2.4)])

    # Alternative callout
    add_solid_rect(s, Inches(8.9), Inches(5.25), Inches(4.1), Inches(1.85), LIGHT_GRAY)
    add_text(s, "Alternative worth considering",
             Inches(9.05), Inches(5.32), Inches(4.0), Inches(0.35),
             font_size=13, bold=True, color=NAVY)
    add_text(s,
             ["Ephemeral QA — spin up Prod-sized QA only during scheduled load-test windows (e.g. 2 days per release), tear down after.",
              "Saves ~80% of QA compute cost in exchange for IaC / automation maturity.",
              "Especially attractive on K8s (kubectl scale --replicas=0) or with Terraform-driven VMs."],
             Inches(9.05), Inches(5.68), Inches(4.0), Inches(1.35),
             font_size=10, color=DARK, anchor=MSO_ANCHOR.TOP)
    return s


def slide_network(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Components to provision (on-prem)",
               "What your network, PKI, and platform teams need to deliver before first install")

    data = [
        ["Component", "Owner", "Purpose"],
        ["F5 BIG-IP DNS (GSLB) with cross-DC health checks, 30s TTL", "Network team", "Cross-DC failover and traffic distribution"],
        ["F5 BIG-IP LTM per DC (HA pair, active/passive)", "Network team", "Per-DC L7 load balancing across Flex Gateway replicas"],
        ["VMs or K8s nodes per DC, sized per slide 6", "Platform team", "Flex Gateway runtime hosts"],
        ["Internal CA + cert lifecycle automation (cert-manager / Venafi)", "PKI team", "mTLS for internal listener + backend, public-CA cert for external listener"],
        ["HashiCorp Vault with performance replication across DCs", "Security platform team", "Secrets, runtime registration token, backend credentials"],
        ["Git repo + CI/CD pipeline (Jenkins / GitHub Actions / Tekton)", "Platform team", "Source of truth for config; pushes to both DCs in parallel"],
        ["Ansible Tower / AWX per DC", "Platform team", "Per-DC config deployment + SIGHUP hot reload (Local mode)"],
        ["Outbound HTTPS to anypoint.mulesoft.com on 443 (Connected mode only)", "Network team", "Policy push + telemetry to Anypoint Control Plane"],
        ["Private DNS — internal zone with backend hostnames", "DNS team", "Backend service resolution (e.g. biztalk-prod.internal.yourco)"],
        ["Perimeter firewall + east-west FW rules", "Network team", "Inbound 443/8443 to gateway; outbound 443 to MS stack and IdP"],
        ["NTP servers (Stratum-2 ideally)", "Sysadmins", "JWT exp/nbf validation requires < 60s skew"],
        ["SIEM (Splunk) + APM (Datadog) ingest endpoints", "Observability", "Dual-shipping from each Flex Gateway replica"],
        ["Backup target for /etc/flex-gateway + Vault snapshots", "Backup team", "Disaster recovery; daily cadence"],
        ["DR-site provisioning (3rd DC or cloud) — optional", "Platform team", "Multi-region geographic resilience beyond 2 metro DCs"],
    ]
    make_table(s, data, Inches(0.3), Inches(1.55), Inches(12.7), Inches(5.45),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=10,
               col_widths=[Inches(5.3), Inches(2.5), Inches(4.9)])
    return s


def slide_dr_postures(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Disaster recovery — posture options",
               "Pick the RTO/RPO point that matches your business continuity requirement and budget")

    data = [
        ["Posture", "Description", "RTO", "RPO", "Cost (vs single DC)"],
        ["A. Cold standby", "DR site provisioned but stopped; manual power-on + DNS flip on failure", "30–60 min", "Zero for config (declarative)", "1.1×"],
        ["B. Warm standby (active/passive)", "DR site running but idle; GSLB sends 100% to primary; on failure flips to DR", "1–5 min", "Zero", "1.8×"],
        ["C. Active/active across DCs in same metro  (recommended)", "Both DCs serve traffic; GSLB weighted round-robin; DC failure = GSLB removes node", "< 60 s", "Zero", "2.0×"],
        ["D. Multi-region active/active", "DCs in different cities/regions (latency > 50ms apart); geo-affinity routing", "< 60 s", "Near-zero (config) / backend-dependent (data)", "2.2×"],
    ]
    make_table(s, data, Inches(0.3), Inches(1.55), Inches(12.7), Inches(2.4),
               header_bg=NAVY, header_fg=WHITE, header_font_size=12, body_font_size=11,
               col_widths=[Inches(3.0), Inches(4.5), Inches(1.4), Inches(2.4), Inches(1.4)])

    # Recommendation callout
    add_solid_rect(s, Inches(0.3), Inches(4.1), Inches(12.7), Inches(1.0), LIGHT_BLUE)
    add_text(s, "Recommendation for citizen-data workload",
             Inches(0.5), Inches(4.18), Inches(8), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    add_text(s,
             "Posture C — active/active across DC-1 and DC-2 within the same metro. RTO < 60s with F5 GSLB + 30s DNS TTL is achievable and proven. Multi-region (D) only when geographic resilience is explicitly required.",
             Inches(0.5), Inches(4.5), Inches(12.4), Inches(0.55),
             font_size=12, color=DARK)

    # DNS failover math
    add_text(s, "DNS failover math (Posture C)",
             Inches(0.5), Inches(5.25), Inches(6), Inches(0.35),
             font_size=13, bold=True, color=NAVY)
    failover = [
        ["Component", "Time"],
        ["GSLB health check interval", "5–15 s"],
        ["Failure detection (3 consecutive misses)", "15–45 s"],
        ["DNS TTL expiry on clients", "≤ 30 s"],
        ["TOTAL worst-case RTO", "≈ 60 s"],
    ]
    make_table(s, failover, Inches(0.5), Inches(5.65), Inches(6.0), Inches(1.5),
               header_bg=BLUE, header_fg=WHITE, header_font_size=11, body_font_size=11,
               col_widths=[Inches(4.0), Inches(2.0)])

    add_text(s, "Note on persistent connections",
             Inches(6.8), Inches(5.25), Inches(6), Inches(0.35),
             font_size=13, bold=True, color=NAVY)
    add_text(s,
             ["Long-lived TLS / gRPC streams won't reconnect on DNS-only failover.",
              "Mitigations: TCP RST from the LB to force reconnect; connection-draining timers on the LTM; short max-connection-age limits set by the gateway. Document this in the DR runbook for affected clients."],
             Inches(6.8), Inches(5.62), Inches(6.2), Inches(1.5),
             font_size=11, color=DARK, anchor=MSO_ANCHOR.TOP)
    return s


def slide_config_replication(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Config replication — Connected vs Local mode",
               "How declarative state stays in lockstep across replicas and DCs in real time")

    # Left half — Connected mode
    add_solid_rect(s, Inches(0.4), Inches(1.55), Inches(6.2), Inches(5.4), LIGHT_BLUE)
    add_text(s, "Connected mode — Anypoint as source of truth",
             Inches(0.55), Inches(1.65), Inches(6.0), Inches(0.4),
             font_size=15, bold=True, color=NAVY)
    add_text(s,
             "Every replica maintains an outbound long-polled HTTPS connection to anypoint.mulesoft.com. Policy changes in API Manager push to all connected replicas globally within seconds. Each replica caches policy locally so brief Anypoint outages don't drop traffic.",
             Inches(0.55), Inches(2.1), Inches(6.0), Inches(1.6),
             font_size=11, color=DARK)

    conn_table = [
        ["Property", "Value"],
        ["Replication mechanism", "Anypoint Control Plane push"],
        ["Push latency (change → live)", "2–10 s end-to-end"],
        ["RPO for config", "Zero"],
        ["Customer infrastructure needed", "None"],
        ["Resilience to Anypoint outage", "Local cache ~24 h"],
        ["Network requirement", "Outbound 443 to anypoint.mulesoft.com"],
    ]
    make_table(s, conn_table, Inches(0.55), Inches(4.0), Inches(5.9), Inches(2.85),
               header_bg=BLUE, header_fg=WHITE, header_font_size=11, body_font_size=10,
               col_widths=[Inches(3.0), Inches(2.9)])

    # Right half — Local mode
    add_solid_rect(s, Inches(6.85), Inches(1.55), Inches(6.2), Inches(5.4), LIGHT_GRAY)
    add_text(s, "Local mode — Git + Ansible AWX",
             Inches(7.0), Inches(1.65), Inches(6.0), Inches(0.4),
             font_size=15, bold=True, color=NAVY)
    add_text(s,
             "Config (declarative YAML) lives in a Git repo — single source of truth. CI pipeline pushes to BOTH DCs in parallel on merge to main. Per-DC Ansible Tower / AWX runs the file deploy + SIGHUP hot reload. Daily reconciliation drift-detects and corrects.",
             Inches(7.0), Inches(2.1), Inches(6.0), Inches(1.6),
             font_size=11, color=DARK)

    local_table = [
        ["Property", "Value"],
        ["Replication mechanism", "GitOps + Ansible AWX (parallel push)"],
        ["Push latency (merge → live)", "< 30 s typical"],
        ["RPO for config", "Zero (Git is authoritative)"],
        ["Customer infrastructure needed", "Git + CI + AWX per DC + Vault for secrets"],
        ["Resilience to GitOps outage", "Existing config keeps serving"],
        ["Network requirement", "None outbound — fully air-gappable"],
    ]
    make_table(s, local_table, Inches(7.0), Inches(4.0), Inches(5.95), Inches(2.85),
               header_bg=BLUE, header_fg=WHITE, header_font_size=11, body_font_size=10,
               col_widths=[Inches(3.0), Inches(2.95)])
    return s


def slide_ops_and_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Operational ownership · risks · next steps",
               "What changes when you take on-prem and how we get to go-live")

    # Ownership shift (left two-thirds)
    add_text(s, "Operational ownership shift (vs SaaS)",
             Inches(0.4), Inches(1.55), Inches(8), Inches(0.35),
             font_size=14, bold=True, color=NAVY)
    own = [
        ["Concern", "SaaS", "On-prem"],
        ["OS / patching", "MuleSoft", "You"],
        ["Runtime upgrades", "MuleSoft", "You"],
        ["HA setup", "Built-in multi-AZ", "Active/active across DCs — you build"],
        ["Capacity planning", "Auto-scale", "You provision for peak"],
        ["Monitoring", "Anypoint built-in + SIEM", "Your SOC + SIEM only"],
        ["Personnel access to runtime", "MuleSoft + you (contractual)", "You only"],
        ["Compliance audit posture", "Inherited (SOC 2, FedRAMP Mod)", "Yours to maintain"],
        ["Release-window downtime", "Rolling, transparent", "Your change window"],
    ]
    make_table(s, own, Inches(0.4), Inches(1.95), Inches(8.4), Inches(4.5),
               header_bg=NAVY, header_fg=WHITE, header_font_size=11, body_font_size=10,
               col_widths=[Inches(2.6), Inches(2.9), Inches(2.9)])

    # Risks (right column top)
    add_text(s, "Key risks", Inches(9.1), Inches(1.55), Inches(4), Inches(0.35),
             font_size=14, bold=True, color=RED_NO)
    risks = [
        "Self-managed = your patching, your incident response, your MTTR",
        "K8s + Flex upgrade collisions if not decoupled",
        "Config drift between DCs without disciplined GitOps",
        "Secrets-sync lag if Vault replication degrades",
        "GSLB itself a SPOF — mitigate with multi-site BIG-IP DNS sync",
    ]
    add_bullets(s, risks, Inches(9.1), Inches(1.95), Inches(4.0), Inches(2.3),
                font_size=10, color=DARK, line_spacing=1.10)

    # Next steps (right column bottom)
    add_text(s, "Next steps", Inches(9.1), Inches(4.4), Inches(4), Inches(0.35),
             font_size=14, bold=True, color=GREEN_OK)
    steps = [
        "Decide VM vs K8s install path with platform team",
        "Network team starts F5 GSLB + LTM change requests",
        "PKI team confirms internal CA + cert-rotation automation",
        "Stand up Vault HA with cross-DC replication first",
        "Build CI pipeline; test config push to a single dev replica",
        "DR drill before go-live — verify < 60 s RTO end-to-end",
    ]
    add_bullets(s, steps, Inches(9.1), Inches(4.8), Inches(4.0), Inches(2.3),
                font_size=10, color=DARK, line_spacing=1.10)
    return s


# -------- Main --------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_why_onprem,
        slide_architecture,
        slide_modes,
        slide_vm_vs_k8s,
        slide_sizing,
        slide_env_matrix,
        slide_network,
        slide_dr_postures,
        slide_config_replication,
        slide_ops_and_next,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        s = build(prs)
        if i != 1:
            add_footer(s, i, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"OK: wrote {OUT_PATH} with {total} slides")


if __name__ == "__main__":
    main()
